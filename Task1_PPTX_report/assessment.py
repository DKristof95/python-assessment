import sys
import json
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Cm
import pandas as pd
from matplotlib import pyplot as plt
import io

try:
    try:
        if sys.argv[1][len(sys.argv[1])-5:len(sys.argv[1])] == ".json":
            f = open(sys.argv[1])
        else:
            f = open(sys.argv[1]+".json")
        config = json.load(f)

        # This makes it easier to create the new slides.
        layout_dict = {
            'title' : 0,
            'text' : 5,
            'list' : 1,
            'picture' : 5,
            'plot' : 5
        }

        prs = Presentation()
        slide_layouts = prs.slide_layouts

        for slide in config['presentation']:
            new_slide = prs.slides.add_slide(slide_layouts[layout_dict[slide['type']]])
            new_slide.shapes.title.text = slide['title']
            left = Cm(3.5)
            top = Cm(3)

            if slide['type'] == 'title':
                new_slide.placeholders[1].text = slide['content']

            elif slide['type'] == 'text':
                width = Cm(14.5)
                height = Cm(30)
                txBox = new_slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = slide['content']

            elif slide['type'] == 'list':
                tf = new_slide.shapes.placeholders[1].text_frame
                for li in slide['content']:
                    p = tf.add_paragraph()
                    p.text = li['text']
                    p.level = li['level']

            elif slide['type'] == 'picture':
                try:
                    new_slide.shapes.add_picture(slide['content'], left, top)
                except FileNotFoundError:
                    print("No picture was found with this name!")

            elif slide['type'] == 'plot':
                try:
                    df = pd.read_csv(slide['content'], sep=';', header=None)

                    plt.plot(df.T.to_numpy()[0], df.T.to_numpy()[1])
                    plt.xlabel(slide['configuration']['x-label'])
                    plt.ylabel(slide['configuration']['y-label'])

                    image_stream = io.BytesIO()
                    plt.savefig(image_stream)
                    new_slide.shapes.add_picture(image_stream, left, top)
                except FileNotFoundError:
                    print("No data file was found with this name!")
                except pd.errors.ParserError:
                    print("Corrupted or badly formatted data file!")

        prs.save('output.pptx')
    except FileNotFoundError:
        print("No configuration file was found with this name!")
    except KeyError:
        print("Corrupted or badly formatted configuration file!")
    except json.decoder.JSONDecodeError:
        print("Corrupted or badly formatted configuration file!")
except IndexError:
    print("No configuration file given!")